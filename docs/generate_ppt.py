"""
Gemicates — Gems Tracker
PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ────────────────────────────────────────────────────────────
BRAND_PURPLE   = RGBColor(0x6C, 0x3F, 0xE8)   # primary
BRAND_INDIGO   = RGBColor(0x43, 0x38, 0xCA)   # accent
BRAND_TEAL     = RGBColor(0x06, 0xB6, 0xD4)   # highlight
BRAND_EMERALD  = RGBColor(0x10, 0xB9, 0x81)   # success
BRAND_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # warning
BRAND_RED      = RGBColor(0xEF, 0x44, 0x44)   # danger
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_PURPLE   = RGBColor(0xED, 0xE9, 0xFE)
DARK_BG        = RGBColor(0x0F, 0x0A, 0x2A)   # dark hero
GRAY_TEXT      = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT      = RGBColor(0x1E, 0x29, 0x3B)
CARD_BG        = RGBColor(0xF8, 0xF7, 0xFF)
BORDER_COLOR   = RGBColor(0xDD, 0xD6, 0xFE)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helper: solid fill ───────────────────────────────────────────────────────
def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

# ── Helper: no fill ──────────────────────────────────────────────────────────
def no_fill(shape):
    shape.fill.background()

# ── Helper: add rectangle ───────────────────────────────────────────────────
def rect(slide, l, t, w, h, color, radius=False):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    solid(shp, color)
    shp.line.fill.background()
    return shp

# ── Helper: text box ─────────────────────────────────────────────────────────
def textbox(slide, text, l, t, w, h,
            size=18, bold=False, color=DARK_TEXT,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

# ── Helper: multi-line textbox ───────────────────────────────────────────────
def multiline(slide, lines, l, t, w, h,
              size=13, bold=False, color=DARK_TEXT,
              align=PP_ALIGN.LEFT, line_spacing=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bd, col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(sz or size)
        run.font.bold  = bd if bd is not None else bold
        run.font.color.rgb = col or color
    return txb

# ── Helper: card ─────────────────────────────────────────────────────────────
def card(slide, l, t, w, h, fill=CARD_BG, border=BORDER_COLOR):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    solid(shp, fill)
    shp.line.color.rgb = border
    shp.line.width     = Pt(0.75)
    return shp

# ── Helper: dot bullet ───────────────────────────────────────────────────────
def bullet_list(slide, items, l, t, w, size=12, color=DARK_TEXT, dot_color=BRAND_PURPLE, gap=0.32):
    for i, item in enumerate(items):
        dy = t + i * gap
        rect(slide, l, dy + 0.08, 0.06, 0.06, dot_color)
        textbox(slide, item, l + 0.12, dy, w - 0.12, 0.3,
                size=size, color=color)

# ── Helper: stat card ────────────────────────────────────────────────────────
def stat_card(slide, l, t, w, h, value, label, val_color=BRAND_PURPLE):
    card(slide, l, t, w, h)
    textbox(slide, value, l + 0.15, t + 0.18, w - 0.3, 0.55,
            size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    textbox(slide, label, l + 0.1, t + 0.72, w - 0.2, 0.35,
            size=10, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ── Helper: section header accent bar ────────────────────────────────────────
def section_bar(slide, label, color=BRAND_PURPLE):
    rect(slide, 0, 0, 13.33, 1.15, DARK_BG)
    rect(slide, 0, 1.15, 13.33, 0.07, color)
    textbox(slide, label, 0.6, 0.2, 12, 0.75,
            size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Helper: icon placeholder circle ──────────────────────────────────────────
def icon_circle(slide, emoji, l, t, r, bg):
    shp = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(r), Inches(r))
    solid(shp, bg)
    shp.line.fill.background()
    textbox(slide, emoji, l, t + 0.02, r, r,
            size=22, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO / TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Background
rect(slide, 0, 0, 13.33, 7.5, DARK_BG)

# Decorative blobs
shp = slide.shapes.add_shape(9, Inches(9.5), Inches(-1.2), Inches(4.5), Inches(4.5))
solid(shp, BRAND_PURPLE)
shp.line.fill.background()
shp.fill.fore_color.theme_color   # no-op, just let solid stand
# set transparency via XML hack
from pptx.oxml.ns import qn
from lxml import etree
shp.fill.fore_color.rgb = RGBColor(0x6C, 0x3F, 0xE8)

shp2 = slide.shapes.add_shape(9, Inches(-1.5), Inches(4.8), Inches(3.5), Inches(3.5))
solid(shp2, BRAND_TEAL)
shp2.line.fill.background()

# Gem icon area (stylised diamond)
gem = slide.shapes.add_shape(1, Inches(0.7), Inches(1.6), Inches(1.1), Inches(1.1))
solid(gem, BRAND_PURPLE)
gem.line.fill.background()
textbox(slide, '💎', 0.7, 1.55, 1.1, 1.1, size=36, align=PP_ALIGN.CENTER)

# Brand name
textbox(slide, 'GEMICATES', 1.95, 1.65, 9, 0.7,
        size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Product name tag
tag = slide.shapes.add_shape(1, Inches(1.95), Inches(2.55), Inches(2.4), Inches(0.38))
solid(tag, BRAND_PURPLE)
tag.line.fill.background()
textbox(slide, 'Gems Tracker', 1.95, 2.55, 2.4, 0.38,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
textbox(slide,
        'The All-in-One Project & Task Intelligence Platform for Modern Teams',
        1.95, 3.1, 8.5, 0.7,
        size=18, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.LEFT)

# Sub-line
textbox(slide,
        'Task Management  ·  Agile Sprints  ·  Time Tracking  ·  OKR Goals  ·  Analytics',
        1.95, 3.9, 8.5, 0.45,
        size=13, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

# Version pill
vpill = slide.shapes.add_shape(1, Inches(1.95), Inches(4.55), Inches(1.3), Inches(0.33))
solid(vpill, RGBColor(0x1E, 0x10, 0x56))
vpill.line.color.rgb = BRAND_PURPLE
vpill.line.width     = Pt(1)
textbox(slide, 'v2.0  |  2026', 1.95, 4.55, 1.3, 0.33,
        size=10, color=RGBColor(0xA5, 0x8A, 0xF8), align=PP_ALIGN.CENTER)

# Bottom band
rect(slide, 0, 6.9, 13.33, 0.6, RGBColor(0x09, 0x06, 0x1A))
textbox(slide, 'Product Feature Presentation  |  Confidential',
        0, 6.92, 13.33, 0.4,
        size=10, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
rect(slide, 0, 1.1,  13.33, 0.06, BRAND_PURPLE)
textbox(slide, "What We'll Cover", 0.6, 0.2, 9, 0.7,
        size=32, bold=True, color=WHITE)

agenda = [
    ('01', 'Product Overview & Vision',         BRAND_PURPLE),
    ('02', 'User Roles & Access Control',        BRAND_INDIGO),
    ('03', 'Core Feature Modules (12)',          BRAND_TEAL),
    ('04', 'Task Management Deep-Dive',          BRAND_EMERALD),
    ('05', 'Agile Sprint Planning',              BRAND_AMBER),
    ('06', 'Goals, Time Tracking & Reports',     BRAND_RED),
    ('07', 'Admin Control Panel',                BRAND_PURPLE),
    ('08', 'Team & Organization Tools',          BRAND_INDIGO),
    ('09', 'Notifications & Announcements',      BRAND_TEAL),
    ('10', 'Technical Architecture',             BRAND_EMERALD),
    ('11', 'Security & Compliance',              BRAND_AMBER),
    ('12', 'Summary & Next Steps',               BRAND_RED),
]

cols = [agenda[:6], agenda[6:]]
for ci, col in enumerate(cols):
    lx = 0.55 + ci * 6.5
    for ri, (num, label, col_color) in enumerate(col):
        ty = 1.35 + ri * 0.97
        numbox = slide.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(0.5), Inches(0.5))
        solid(numbox, col_color)
        numbox.line.fill.background()
        textbox(slide, num, lx, ty + 0.05, 0.5, 0.4,
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(slide, label, lx + 0.62, ty + 0.08, 5.5, 0.38,
                size=13, color=DARK_TEXT)
        rect(slide, lx + 0.62, ty + 0.49, 5.2, 0.015,
             RGBColor(0xE2, 0xE8, 0xF0))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PRODUCT OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '01  Product Overview & Vision', BRAND_PURPLE)

# Left panel — vision text
textbox(slide, 'What is Gemicates?', 0.55, 1.45, 6, 0.45,
        size=20, bold=True, color=BRAND_PURPLE)
textbox(slide,
        'Gemicates (Gems Tracker) is an all-in-one project and task '
        'management platform built for modern teams — combining the '
        'simplicity of task tracking with the power of Agile sprint planning, '
        'OKR goal management, organizational hierarchy, and deep analytics.',
        0.55, 1.95, 5.8, 1.5,
        size=12.5, color=DARK_TEXT)

# Vision pillars
pillars = [
    ('🎯', 'Clarity',        'Every member sees exactly what they need to do — today, this week, this sprint.', BRAND_PURPLE),
    ('📊', 'Accountability', 'Audit logs, time tracking & activity histories make work visible & measurable.', BRAND_INDIGO),
    ('⚡', 'Agility',        'Native sprint boards, backlog grooming & kanban support Scrum out of the box.', BRAND_TEAL),
    ('🔗', 'Alignment',      'OKR Goals connect daily task work to company-level objectives seamlessly.', BRAND_EMERALD),
]
for i, (ico, title, desc, col) in enumerate(pillars):
    ty = 3.65 + i * 0.82
    circle = slide.shapes.add_shape(9, Inches(0.55), Inches(ty - 0.03), Inches(0.45), Inches(0.45))
    solid(circle, col)
    circle.line.fill.background()
    textbox(slide, ico, 0.55, ty - 0.04, 0.45, 0.44, size=14, align=PP_ALIGN.CENTER)
    textbox(slide, title, 1.1, ty - 0.02, 1.6, 0.28, size=12, bold=True, color=col)
    textbox(slide, desc,  2.75, ty - 0.02, 3.4, 0.5, size=10.5, color=GRAY_TEXT)

# Right panel — stats grid
stats = [
    ('150+', 'REST API Endpoints'),
    ('12',   'Core Feature Modules'),
    ('25+',  'Database Models'),
    ('23',   'Frontend Pages'),
    ('3',    'User Role Tiers'),
    ('4',    'Export Formats'),
]
for i, (val, lbl) in enumerate(stats):
    row, col_i = divmod(i, 2)
    sx = 7.0 + col_i * 3.1
    sy = 1.45 + row * 1.95
    stat_card(slide, sx, sy, 2.95, 1.7, val, lbl,
              [BRAND_PURPLE, BRAND_INDIGO, BRAND_TEAL, BRAND_EMERALD, BRAND_AMBER, BRAND_RED][i])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — USER ROLES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '02  User Roles & Access Control', BRAND_INDIGO)

roles = [
    {
        'icon':  '🛡️',
        'title': 'Administrator',
        'color':  BRAND_PURPLE,
        'bg':     RGBColor(0xED, 0xE9, 0xFE),
        'items':  [
            'Full system access',
            'Manage users, entities, projects',
            'Configure site & mail settings',
            'View all reports & audit logs',
            'Create announcements',
            'Manage org chart & departments',
            'Review leave requests',
            'Spawn recurring tasks',
        ],
    },
    {
        'icon':  '👔',
        'title': 'Project Manager',
        'color':  BRAND_TEAL,
        'bg':     RGBColor(0xCC, 0xFB, 0xF1),
        'items':  [
            'Create & manage sprints',
            'Add / remove project members',
            'Assign tasks to team members',
            'Start & complete sprints',
            'Edit & delete sprint issues',
            'Create project milestones',
            'View project-level reports',
            'All member capabilities',
        ],
    },
    {
        'icon':  '👤',
        'title': 'Team Member',
        'color':  BRAND_AMBER,
        'bg':     RGBColor(0xFE, 0xF3, 0xC7),
        'items':  [
            'Create, edit & delete own tasks',
            'Log time & use active timer',
            'Add comments & subtasks',
            'View assigned projects',
            'Participate in sprint boards',
            'Manage personal notes',
            'Save filter presets',
            'Submit leave requests',
        ],
    },
]

for i, role in enumerate(roles):
    lx = 0.55 + i * 4.25
    bx = slide.shapes.add_shape(1, Inches(lx), Inches(1.35), Inches(3.95), Inches(5.85))
    solid(bx, role['bg'])
    bx.line.color.rgb = role['color']
    bx.line.width     = Pt(1.5)

    # Header
    hdr = slide.shapes.add_shape(1, Inches(lx), Inches(1.35), Inches(3.95), Inches(0.9))
    solid(hdr, role['color'])
    hdr.line.fill.background()
    textbox(slide, role['icon'] + '  ' + role['title'],
            lx + 0.15, 1.45, 3.65, 0.65,
            size=15, bold=True, color=WHITE)

    bullet_list(slide, role['items'],
                lx + 0.2, 2.42, 3.55,
                size=11, color=DARK_TEXT, dot_color=role['color'], gap=0.56)

# Access summary note
rect(slide, 0.55, 7.1, 12.2, 0.28, LIGHT_PURPLE)
textbox(slide,
        'Roles are enforced at both API middleware and frontend rendering levels — users only see what they are authorised to access.',
        0.65, 7.11, 12, 0.25, size=9.5, color=BRAND_INDIGO)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE MODULES OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '03  Core Feature Modules', BRAND_TEAL)

modules = [
    ('✅', 'Task Management',           BRAND_PURPLE),
    ('📋', 'Project Management',        BRAND_INDIGO),
    ('⚡', 'Sprint & Agile Planning',   BRAND_TEAL),
    ('⏱️', 'Time Tracking',             BRAND_EMERALD),
    ('🎯', 'Goals & OKRs',              BRAND_AMBER),
    ('📊', 'Reporting & Analytics',     BRAND_RED),
    ('🏢', 'Team & Organization',       BRAND_PURPLE),
    ('🔔', 'Notifications',             BRAND_INDIGO),
    ('📢', 'Announcements',             BRAND_TEAL),
    ('⚙️', 'Admin Control Panel',       BRAND_EMERALD),
    ('📝', 'Personal Productivity',     BRAND_AMBER),
    ('🔐', 'Security & Audit',          BRAND_RED),
]

for i, (ico, name, col) in enumerate(modules):
    row, ci = divmod(i, 4)
    mx = 0.55 + ci * 3.2
    my = 1.4  + row * 2.0
    bx = slide.shapes.add_shape(1, Inches(mx), Inches(my), Inches(3.0), Inches(1.7))
    solid(bx, CARD_BG)
    bx.line.color.rgb = col
    bx.line.width     = Pt(1.2)
    c = slide.shapes.add_shape(9, Inches(mx + 0.15), Inches(my + 0.18), Inches(0.55), Inches(0.55))
    solid(c, col)
    c.line.fill.background()
    textbox(slide, ico, mx + 0.15, my + 0.17, 0.55, 0.55, size=16, align=PP_ALIGN.CENTER)
    textbox(slide, name, mx + 0.82, my + 0.25, 2.05, 0.5, size=11.5, bold=True, color=col)
    # bar accent bottom
    rect(slide, mx, my + 1.5, 3.0, 0.2, col)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TASK MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '04  Task Management', BRAND_PURPLE)

# Left col — fields
textbox(slide, 'Task Fields', 0.55, 1.4, 4, 0.38, size=15, bold=True, color=BRAND_PURPLE)
fields = [
    ('Title',        'Required task name'),
    ('Issue Type',   'Story  ·  Bug  ·  Task  ·  Epic'),
    ('Status',       'Configurable custom statuses'),
    ('Priority',     'Low  ·  Medium  ·  High  ·  Critical'),
    ('Hours Spent',  '0–24 hours manual log'),
    ('Due Date',     'Optional deadline'),
    ('Assignee',     'Team member responsible'),
    ('Story Points', 'Agile effort estimation'),
    ('Recurrence',   'Daily or weekly repeat'),
]
for i, (k, v) in enumerate(fields):
    ty = 1.88 + i * 0.52
    rect(slide, 0.55, ty + 0.07, 0.05, 0.28, BRAND_PURPLE)
    textbox(slide, k + ':',  0.72, ty, 1.5, 0.38, size=11, bold=True, color=DARK_TEXT)
    textbox(slide, v,        2.28, ty, 3.5, 0.38, size=11, color=GRAY_TEXT)

# Middle col — operations & views
textbox(slide, 'Operations', 6.2, 1.4, 3.2, 0.38, size=15, bold=True, color=BRAND_INDIGO)
ops = ['Create / Edit / Delete', 'Soft delete + Restore', 'Bulk Import (Excel .xlsx)',
       'Bulk Delete (multi-select)', 'Export — Excel & PDF', 'Cycle status on click',
       'Reassign to any member', 'Clone recurring tasks']
bullet_list(slide, ops, 6.2, 1.88, 3.2, size=11, dot_color=BRAND_INDIGO)

textbox(slide, 'Views', 6.2, 4.55, 3.2, 0.38, size=15, bold=True, color=BRAND_INDIGO)
views = ['📋 List — sortable, paginated table', '🗂  Kanban — drag-and-drop status columns']
bullet_list(slide, views, 6.2, 5.02, 3.2, size=11, dot_color=BRAND_INDIGO)

# Right col — rich detail features
textbox(slide, 'Task Detail Features', 9.6, 1.4, 3.5, 0.38, size=15, bold=True, color=BRAND_TEAL)
detail_features = [
    ('💬', 'Threaded Comments'),
    ('☑️', 'Subtask Checklist'),
    ('⏱️', 'Active Timer'),
    ('🔗', 'Task Dependencies'),
    ('👁', 'Watchers & Alerts'),
    ('📜', 'Activity Log'),
    ('🏷️', 'Labels & Tags'),
    ('📎', 'Link Types (blocks, depends)'),
]
for i, (ico, lbl) in enumerate(detail_features):
    ty = 1.88 + i * 0.56
    c = slide.shapes.add_shape(9, Inches(9.6), Inches(ty), Inches(0.38), Inches(0.38))
    solid(c, RGBColor(0xEE, 0xE0, 0xFF))
    c.line.fill.background()
    textbox(slide, ico, 9.6, ty, 0.38, 0.38, size=13, align=PP_ALIGN.CENTER)
    textbox(slide, lbl, 10.1, ty + 0.04, 2.95, 0.34, size=11, color=DARK_TEXT)

# Filter strip
rect(slide, 0.55, 7.1, 12.2, 0.28, RGBColor(0xF3, 0xF0, 0xFF))
textbox(slide, 'Filters:  Project  ·  Status  ·  Priority  ·  Date Range  ·  Free-text Search  ·  Saved Presets  ·  Quick: Today / This Week / This Month',
        0.7, 7.11, 12, 0.25, size=9.5, color=BRAND_INDIGO)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AGILE SPRINT PLANNING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '05  Agile Sprint Planning', BRAND_AMBER)

# Sprint lifecycle
textbox(slide, 'Sprint Lifecycle', 0.55, 1.38, 5, 0.38, size=15, bold=True, color=BRAND_AMBER)
phases = [('Planning', BRAND_INDIGO), ('Active', BRAND_EMERALD), ('Completed', GRAY_TEXT)]
for i, (ph, col) in enumerate(phases):
    px = 0.65 + i * 2.55
    b  = slide.shapes.add_shape(1, Inches(px), Inches(1.88), Inches(2.1), Inches(0.55))
    solid(b, col)
    b.line.fill.background()
    textbox(slide, ph, px, 1.88, 2.1, 0.55,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        textbox(slide, '►', px + 2.12, 1.97, 0.35, 0.35,
                size=14, bold=True, color=BRAND_AMBER, align=PP_ALIGN.CENTER)

# Kanban board visual
textbox(slide, 'Kanban Board', 0.55, 2.65, 5, 0.38, size=15, bold=True, color=BRAND_AMBER)
cols_demo = [
    ('To Do',       BRAND_INDIGO,   ['User auth flow', 'API gateway setup']),
    ('In Progress', BRAND_AMBER,    ['Sprint module UI', 'DB migrations']),
    ('Done',        BRAND_EMERALD,  ['Login page', 'JWT integration']),
]
for ci, (col_name, col_color, tasks_demo) in enumerate(cols_demo):
    cx = 0.55 + ci * 2.38
    hdr = slide.shapes.add_shape(1, Inches(cx), Inches(3.12), Inches(2.2), Inches(0.42))
    solid(hdr, col_color)
    hdr.line.fill.background()
    textbox(slide, col_name, cx, 3.12, 2.2, 0.42,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ti, t in enumerate(tasks_demo):
        tb = slide.shapes.add_shape(1, Inches(cx), Inches(3.62 + ti * 0.72), Inches(2.2), Inches(0.6))
        solid(tb, WHITE)
        tb.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        tb.line.width     = Pt(0.75)
        textbox(slide, t, cx + 0.1, 3.67 + ti * 0.72, 2.0, 0.5, size=10, color=DARK_TEXT)

# Right col — features
textbox(slide, 'Sprint Features', 7.3, 1.38, 5.5, 0.38, size=15, bold=True, color=BRAND_AMBER)
sprint_features = [
    '⚡  Drag-and-drop between status columns',
    '🎯  Configurable status columns (custom statuses)',
    '📌  Sprint goal statement',
    '📊  Story point total vs. done progress bar',
    '🗂  Backlog for unspinted issues',
    '👤  Assignee picker on issue create',
    '✏️  Edit & delete issues (managers only)',
    '🔢  Issue type breakdown (Story / Bug / Task / Epic)',
    '🚀  Start sprint & Complete sprint actions',
    '📋  Add issues from backlog to any sprint',
    '🔄  Optimistic UI — instant card movement',
]
bullet_list(slide, sprint_features, 7.3, 1.88, 5.7,
            size=11.5, dot_color=BRAND_AMBER, gap=0.44)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GOALS, TIME TRACKING & REPORTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '06  Goals · Time Tracking · Reports', BRAND_EMERALD)

# Goals / OKR
card(slide, 0.55, 1.38, 3.85, 5.9)
textbox(slide, '🎯  Goals & OKRs', 0.75, 1.52, 3.5, 0.4,
        size=14, bold=True, color=BRAND_EMERALD)
okr_items = [
    'Create objectives with owner & due date',
    'Attach multiple Key Results per goal',
    'Track Current vs. Target values',
    'Custom unit (%, users, revenue...)',
    'Auto-calculated progress %',
    'Goal progress = avg of KR progress',
    'Status: Active / Completed',
    'Expand goal to edit KRs inline',
]
bullet_list(slide, okr_items, 0.75, 2.02, 3.45,
            size=11, dot_color=BRAND_EMERALD, gap=0.5)

# Time Tracking
card(slide, 4.6, 1.38, 3.85, 5.9)
textbox(slide, '⏱️  Time Tracking', 4.8, 1.52, 3.5, 0.4,
        size=14, bold=True, color=BRAND_TEAL)
time_items = [
    'Manual hours per task (0–24)',
    'One active timer per user',
    'Start / Stop with one click',
    'Auto-records seconds elapsed',
    'Time entry history in task detail',
    'Persists across navigation',
    'Weekly summary: this vs. last week',
    'Hours-by-project breakdown',
]
bullet_list(slide, time_items, 4.8, 2.02, 3.45,
            size=11, dot_color=BRAND_TEAL, gap=0.5)

# Reports
card(slide, 8.65, 1.38, 4.15, 5.9)
textbox(slide, '📊  Reports & Analytics', 8.85, 1.52, 3.8, 0.4,
        size=14, bold=True, color=BRAND_INDIGO)
report_sections = [
    ('Admin Reports', [
        'All-tasks report (any user)',
        'User activity heatmap',
        'Tasks per project chart',
        'Hours by project',
        'Deadline alerts dashboard',
        'CSV export',
    ]),
    ('User Reports', [
        'Personal task summary',
        'Status breakdown counts',
        'Hours by project (own)',
        'Weekly task comparison',
        'Activity streak counter',
    ]),
]
ty = 2.02
for section, items in report_sections:
    textbox(slide, section, 8.85, ty, 3.8, 0.34,
            size=11, bold=True, color=BRAND_INDIGO)
    ty += 0.36
    bullet_list(slide, items, 8.85, ty, 3.65,
                size=10.5, dot_color=BRAND_INDIGO, gap=0.44)
    ty += len(items) * 0.44 + 0.25

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ADMIN CONTROL PANEL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '07  Admin Control Panel', BRAND_RED)

admin_cards = [
    {
        'icon': '👥', 'title': 'User Management', 'color': BRAND_PURPLE,
        'items': ['Create users with invite email', 'Edit name, email, role, avatar',
                  'Soft delete & restore', 'Assign to projects & entities',
                  'Toggle active / inactive'],
    },
    {
        'icon': '🏢', 'title': 'Entity & Project Mgmt', 'color': BRAND_TEAL,
        'items': ['Create / edit / delete entities', 'Create projects under entities',
                  'Set project status & dates', 'Soft delete with restore',
                  'View project task stats'],
    },
    {
        'icon': '⚙️', 'title': 'Site Settings', 'color': BRAND_AMBER,
        'items': ['Set site name & upload logo', 'Configure custom task statuses',
                  'Statuses reflect in Kanban & forms', 'Remove logo option',
                  'Public settings API for SPA'],
    },
    {
        'icon': '📧', 'title': 'Mail Settings (SMTP)', 'color': BRAND_RED,
        'items': ['SMTP host, port, auth config', 'Sender name & from address',
                  'App URL setting', 'Password masked in responses',
                  'Send test email to self'],
    },
    {
        'icon': '📜', 'title': 'Audit Log', 'color': BRAND_INDIGO,
        'items': ['Entity-level change history', 'Action: create / update / delete',
                  'Records: who, what, when', 'Old vs. new value diff',
                  'Filter by entity type & ID'],
    },
    {
        'icon': '🌳', 'title': 'Org Chart', 'color': BRAND_EMERALD,
        'items': ['Visual manager-subordinate tree', 'Set manager per user',
                  'Recursive hierarchy display', 'Department management',
                  'Leave request review'],
    },
]

for i, ac in enumerate(admin_cards):
    row, ci = divmod(i, 3)
    ax = 0.55 + ci * 4.25
    ay = 1.38 + row * 2.9
    bx = slide.shapes.add_shape(1, Inches(ax), Inches(ay), Inches(4.0), Inches(2.65))
    solid(bx, CARD_BG)
    bx.line.color.rgb = ac['color']
    bx.line.width     = Pt(1.2)
    hdr = slide.shapes.add_shape(1, Inches(ax), Inches(ay), Inches(4.0), Inches(0.58))
    solid(hdr, ac['color'])
    hdr.line.fill.background()
    textbox(slide, ac['icon'] + '  ' + ac['title'],
            ax + 0.12, ay + 0.1, 3.75, 0.42,
            size=12, bold=True, color=WHITE)
    bullet_list(slide, ac['items'], ax + 0.15, ay + 0.7, 3.7,
                size=10.5, dot_color=ac['color'], gap=0.39)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM & ORGANIZATION + NOTIFICATIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '08 & 09  Team, Organization & Notifications', BRAND_INDIGO)

# Left — Team & Org
card(slide, 0.55, 1.38, 6.0, 5.9)
textbox(slide, '🏢  Team & Organization', 0.75, 1.52, 5.6, 0.4,
        size=15, bold=True, color=BRAND_INDIGO)

org_sections = [
    ('Entities',          ['Top-level business divisions', 'Projects belong to entities', 'Users assigned to entities']),
    ('Project Members',   ['Manager & Member roles', 'Managers create sprints & assign tasks', 'Members view & contribute']),
    ('Departments',       ['Hierarchical parent-child structure', 'Department head assignment', 'Nested org chart']),
    ('Leave Management',  ['Types: Annual · Sick · Personal', 'Request with date range & reason', 'Admin: Approve or Reject']),
]
ty = 2.0
for title, items in org_sections:
    rect(slide, 0.75, ty, 0.04, len(items) * 0.42 + 0.06, BRAND_INDIGO)
    textbox(slide, title, 0.92, ty - 0.04, 4.5, 0.34, size=12, bold=True, color=BRAND_INDIGO)
    for item in items:
        textbox(slide, '·  ' + item, 0.92, ty + 0.32, 5.2, 0.36, size=10.5, color=DARK_TEXT)
        ty += 0.42
    ty += 0.35

# Right — Notifications
card(slide, 6.78, 1.38, 6.0, 2.9)
textbox(slide, '🔔  Notifications', 6.98, 1.52, 5.6, 0.4,
        size=15, bold=True, color=BRAND_TEAL)
notif_triggers = [
    ('Task assigned to you',        'You have been assigned: [Task Title]'),
    ('Added to a project',          'You have been added to: [Project]'),
    ('New sprint started',          'Sprint [Name] is now active'),
    ('New announcement published',  'New announcement: [Title]'),
]
for i, (trigger, msg) in enumerate(notif_triggers):
    ty_n = 2.05 + i * 0.56
    rect(slide, 6.98, ty_n + 0.1, 0.18, 0.18, BRAND_TEAL)
    textbox(slide, trigger, 7.28, ty_n, 3.2, 0.28, size=10.5, bold=True, color=DARK_TEXT)
    textbox(slide, '→ ' + msg, 7.28, ty_n + 0.26, 5.2, 0.26, size=10, color=GRAY_TEXT, italic=True)

# Announcements
card(slide, 6.78, 4.48, 6.0, 2.82)
textbox(slide, '📢  Announcements', 6.98, 4.62, 5.6, 0.4,
        size=15, bold=True, color=BRAND_AMBER)
ann_items = [
    'Admins create announcements with title & body',
    'All active users notified instantly on publish',
    'Per-user read / unread tracking',
    'Announcements page with read status badge',
    'Admins can delete announcements',
    'Bell icon shows unread count badge',
    'Mark individual or all-read actions',
]
bullet_list(slide, ann_items, 6.98, 5.12, 5.5,
            size=11, dot_color=BRAND_AMBER, gap=0.46)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PERSONAL PRODUCTIVITY TOOLS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '10  Personal Productivity Tools', BRAND_EMERALD)

tools = [
    {
        'icon': '📝', 'title': 'Notes',
        'color': BRAND_PURPLE,
        'desc': 'Private, per-user notes with pin support for quick access.',
        'items': ['Title + rich body', 'Pin to top', 'Create / Edit / Delete'],
    },
    {
        'icon': '🔖', 'title': 'Filter Presets',
        'color': BRAND_TEAL,
        'desc': 'Save complex filter combinations and load them in one click.',
        'items': ['Name and save any filter', 'Instant reload', 'Delete when done'],
    },
    {
        'icon': '📌', 'title': 'Pinned Projects',
        'color': BRAND_AMBER,
        'desc': 'Pin your most-used projects for instant dashboard access.',
        'items': ['Quick-pin / unpin', 'Progress bar per project', 'Shown on dashboard'],
    },
    {
        'icon': '🔥', 'title': 'Activity Streak',
        'color': BRAND_RED,
        'desc': 'Consecutive days with at least one task logged — stay on track.',
        'items': ['Visible on user dashboard', 'Motivates daily logging', 'Auto-calculated'],
    },
    {
        'icon': '⌨️', 'title': 'Keyboard Shortcuts',
        'color': BRAND_INDIGO,
        'desc': 'Power-user speed without leaving the keyboard.',
        'items': ['T — New task', 'Ctrl+N — New task', 'More shortcuts planned'],
    },
    {
        'icon': '📅', 'title': 'Dashboard Calendar',
        'color': BRAND_EMERALD,
        'desc': 'Mini calendar on the user dashboard for at-a-glance scheduling.',
        'items': ['Today highlighted', 'Task counts per day', 'Quick date navigation'],
    },
]

for i, tool in enumerate(tools):
    row, ci = divmod(i, 3)
    tx = 0.55 + ci * 4.25
    ty = 1.38 + row * 2.88
    bx = slide.shapes.add_shape(1, Inches(tx), Inches(ty), Inches(4.0), Inches(2.65))
    solid(bx, CARD_BG)
    bx.line.color.rgb = tool['color']
    bx.line.width     = Pt(1)
    c = slide.shapes.add_shape(9, Inches(tx + 0.15), Inches(ty + 0.15), Inches(0.5), Inches(0.5))
    solid(c, tool['color'])
    c.line.fill.background()
    textbox(slide, tool['icon'], tx + 0.15, ty + 0.14, 0.5, 0.5,
            size=16, align=PP_ALIGN.CENTER)
    textbox(slide, tool['title'], tx + 0.78, ty + 0.2, 3.0, 0.36,
            size=13, bold=True, color=tool['color'])
    textbox(slide, tool['desc'],  tx + 0.15, ty + 0.75, 3.7, 0.5,
            size=10, color=GRAY_TEXT)
    bullet_list(slide, tool['items'], tx + 0.15, ty + 1.3, 3.7,
                size=10.5, dot_color=tool['color'], gap=0.4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '11  Technical Architecture', BRAND_TEAL)

# Stack table left
textbox(slide, 'Technology Stack', 0.55, 1.38, 4.5, 0.4,
        size=15, bold=True, color=BRAND_TEAL)
stack = [
    ('Frontend',         'React 18 + TypeScript, Vite'),
    ('State / Cache',    'TanStack Query v5, Zustand'),
    ('UI',               'Tailwind CSS v3'),
    ('Forms',            'React Hook Form + Zod'),
    ('Drag & Drop',      '@dnd-kit/core'),
    ('Backend',          'PHP (Custom MVC Framework)'),
    ('Database',         'PostgreSQL'),
    ('Auth',             'JWT — 24-hour tokens'),
    ('Mail',             'SMTP (admin-configurable)'),
    ('File Storage',     'Local server storage'),
]
for i, (layer, tech) in enumerate(stack):
    ty = 1.88 + i * 0.5
    rect(slide, 0.55, ty + 0.08, 0.05, 0.28, BRAND_TEAL)
    textbox(slide, layer + ':', 0.72, ty, 1.7, 0.38, size=11, bold=True, color=DARK_TEXT)
    textbox(slide, tech,       2.45, ty, 3.3, 0.38, size=11, color=GRAY_TEXT)

# Architecture diagram (text-based)
textbox(slide, 'Architecture Flow', 6.2, 1.38, 7.0, 0.4,
        size=15, bold=True, color=BRAND_TEAL)

arch_layers = [
    ('🌐  Browser', 'React SPA — 23 pages, 25 API clients', BRAND_PURPLE),
    ('🔀  REST API', 'PHP Router — 150+ endpoints', BRAND_INDIGO),
    ('🧠  Controllers', '30 controllers — Business logic layer', BRAND_TEAL),
    ('📦  Models', '25+ Core\\Model classes (PDO/PostgreSQL)', BRAND_EMERALD),
    ('🗄️  Database', 'PostgreSQL — Soft deletes, Audit log', BRAND_AMBER),
]
for i, (layer, desc, col) in enumerate(arch_layers):
    ay = 1.88 + i * 1.02
    bx = slide.shapes.add_shape(1, Inches(6.2), Inches(ay), Inches(6.8), Inches(0.62))
    solid(bx, CARD_BG)
    bx.line.color.rgb = col
    bx.line.width     = Pt(1.2)
    accent = slide.shapes.add_shape(1, Inches(6.2), Inches(ay), Inches(0.18), Inches(0.62))
    solid(accent, col)
    accent.line.fill.background()
    textbox(slide, layer, 6.48, ay + 0.1, 2.8, 0.38, size=11.5, bold=True, color=col)
    textbox(slide, desc,  9.35, ay + 0.1, 3.55, 0.38, size=11, color=GRAY_TEXT)
    if i < 4:
        textbox(slide, '▼', 9.5, ay + 0.65, 0.5, 0.35,
                size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SECURITY & COMPLIANCE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '12  Security & Compliance', BRAND_RED)

security_panels = [
    {
        'icon': '🔐', 'title': 'Authentication', 'color': BRAND_PURPLE,
        'items': ['JWT tokens — 24-hour expiry', 'Bcrypt password hashing (cost 12)',
                  'No plaintext credentials stored', 'Token validated on every API call'],
    },
    {
        'icon': '🛡️', 'title': 'Authorization', 'color': BRAND_INDIGO,
        'items': ['Role-based middleware on every route', 'Users access only own data',
                  'Admin-only routes protected', 'Manager actions scoped to project'],
    },
    {
        'icon': '💾', 'title': 'Data Protection', 'color': BRAND_TEAL,
        'items': ['Soft deletes — all records restorable', 'Entity-level audit log',
                  'File uploads: MIME + 2MB limit', 'No permanent deletions'],
    },
    {
        'icon': '🔍', 'title': 'Input Validation', 'color': BRAND_EMERALD,
        'items': ['Server-side sanitization (PHP)', 'PDO prepared statements (SQL injection)',
                  'Zod schema validation (frontend)', 'Context-aware output encoding (XSS)'],
    },
    {
        'icon': '📬', 'title': 'Mail Security', 'color': BRAND_AMBER,
        'items': ['SMTP credentials in .env file', 'Password masked in API responses',
                  'Test email before go-live', 'Admin-only mail settings access'],
    },
    {
        'icon': '📋', 'title': 'Audit & Compliance', 'color': BRAND_RED,
        'items': ['Full audit trail per entity', 'Who, what, when recorded',
                  'Old vs. new value diff stored', 'Filterable by type & entity ID'],
    },
]

for i, panel in enumerate(security_panels):
    row, ci = divmod(i, 3)
    sx = 0.55 + ci * 4.25
    sy = 1.38 + row * 2.88
    bx = slide.shapes.add_shape(1, Inches(sx), Inches(sy), Inches(4.0), Inches(2.65))
    solid(bx, CARD_BG)
    bx.line.color.rgb = panel['color']
    bx.line.width     = Pt(1)
    hdr = slide.shapes.add_shape(1, Inches(sx), Inches(sy), Inches(4.0), Inches(0.6))
    solid(hdr, panel['color'])
    hdr.line.fill.background()
    textbox(slide, panel['icon'] + '  ' + panel['title'],
            sx + 0.12, sy + 0.1, 3.75, 0.44,
            size=12, bold=True, color=WHITE)
    bullet_list(slide, panel['items'], sx + 0.15, sy + 0.72, 3.7,
                size=10.5, dot_color=panel['color'], gap=0.46)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — FEATURE MATRIX
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xFD, 0xFC, 0xFF))
section_bar(slide, '13  Feature Matrix by Role', BRAND_PURPLE)

matrix_header = ['Feature', 'Admin', 'Manager', 'Member']
matrix_rows = [
    # TASKS
    ('─── TASKS ───────────────────────', '',    '',    ''),
    ('Create / Edit / Delete own tasks', '✅', '✅', '✅'),
    ('Bulk import (Excel)',              '✅', '✅', '✅'),
    ('Export tasks (Excel / PDF)',       '✅', '✅', '✅'),
    ('Assign tasks to others',           '✅', '✅', '❌'),
    ('Edit or delete ANY task',          '✅', '❌', '❌'),
    ('Spawn recurring tasks',            '✅', '❌', '❌'),
    # SPRINTS
    ('─── SPRINTS ─────────────────────', '',    '',    ''),
    ('Create / edit / delete sprints',   '✅', '✅', '❌'),
    ('Start & complete sprints',         '✅', '✅', '❌'),
    ('Edit & delete sprint issues',      '✅', '✅', '❌'),
    ('View & drag sprint board',         '✅', '✅', '✅'),
    # ADMIN
    ('─── ADMIN ───────────────────────', '',    '',    ''),
    ('Manage users',                     '✅', '❌', '❌'),
    ('Site & mail settings',             '✅', '❌', '❌'),
    ('Create announcements',             '✅', '❌', '❌'),
    ('View audit log',                   '✅', '❌', '❌'),
    # PERSONAL
    ('─── PERSONAL ────────────────────', '',    '',    ''),
    ('Notes, Presets, Pinned Projects',  '✅', '✅', '✅'),
    ('Submit leave requests',            '✅', '✅', '✅'),
    ('OKR Goals',                        '✅', '✅', '✅'),
    ('Admin-wide reports',               '✅', '❌', '❌'),
]

col_colors = [BRAND_PURPLE, BRAND_TEAL, BRAND_AMBER]
col_xs     = [5.8, 8.5, 11.2]

# Header row
for ci, (ch, cx) in enumerate(zip(matrix_header[1:], col_xs)):
    hbox = slide.shapes.add_shape(1, Inches(cx), Inches(1.38), Inches(2.0), Inches(0.48))
    solid(hbox, col_colors[ci])
    hbox.line.fill.background()
    textbox(slide, ch, cx, 1.38, 2.0, 0.48,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, 'Feature', 0.55, 1.38, 5.0, 0.48,
        size=12, bold=True, color=DARK_TEXT)

for ri, row in enumerate(matrix_rows):
    ry = 1.95 + ri * 0.24
    bg = RGBColor(0xF8, 0xF7, 0xFF) if ri % 2 == 0 else WHITE
    is_section = row[0].startswith('───')
    if is_section:
        bg = RGBColor(0xED, 0xE9, 0xFE)
    bx = slide.shapes.add_shape(1, Inches(0.45), Inches(ry), Inches(12.45), Inches(0.24))
    solid(bx, bg)
    bx.line.fill.background()
    textbox(slide, row[0], 0.55, ry, 5.0, 0.24,
            size=9.5,
            bold=is_section,
            color=BRAND_INDIGO if is_section else DARK_TEXT)
    for ci, (val, cx) in enumerate(zip(row[1:], col_xs)):
        c = BRAND_EMERALD if val == '✅' else (BRAND_RED if val == '❌' else GRAY_TEXT)
        textbox(slide, val, cx, ry, 2.0, 0.24,
                size=10, color=c, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING / THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, DARK_BG)

# Top glow blob
shp = slide.shapes.add_shape(9, Inches(4.5), Inches(-2), Inches(6), Inches(6))
solid(shp, BRAND_PURPLE)
shp.line.fill.background()

shp2 = slide.shapes.add_shape(9, Inches(-1), Inches(3.5), Inches(4), Inches(4))
solid(shp2, BRAND_TEAL)
shp2.line.fill.background()

# Gem icon
textbox(slide, '💎', 5.7, 0.9, 1.6, 1.2, size=52, align=PP_ALIGN.CENTER)

# Title
textbox(slide, 'GEMICATES', 0, 2.1, 13.33, 0.8,
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(slide, 'Gems Tracker  —  The Future of Team Productivity', 0, 2.95, 13.33, 0.5,
        size=17, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

# Divider
rect(slide, 3.5, 3.65, 6.33, 0.04, BRAND_PURPLE)

# Summary stats row
summary = [('150+', 'API Endpoints'), ('12', 'Modules'), ('25+', 'Data Models'), ('3', 'Role Tiers')]
for i, (v, l) in enumerate(summary):
    sx = 2.2 + i * 2.35
    textbox(slide, v, sx, 3.85, 2.2, 0.55,
            size=32, bold=True, color=BRAND_TEAL, align=PP_ALIGN.CENTER)
    textbox(slide, l, sx, 4.4, 2.2, 0.35,
            size=11, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

# CTA pill
pill = slide.shapes.add_shape(1, Inches(4.5), Inches(5.1), Inches(4.33), Inches(0.62))
solid(pill, BRAND_PURPLE)
pill.line.fill.background()
textbox(slide, '🚀  Ready to Get Started with Gemicates?', 4.5, 5.1, 4.33, 0.62,
        size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(slide, 'Task Management  ·  Agile Sprints  ·  Time Tracking  ·  OKRs  ·  Analytics  ·  Team Management',
        0, 5.95, 13.33, 0.4,
        size=11, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

# Footer
rect(slide, 0, 6.9, 13.33, 0.6, RGBColor(0x09, 0x06, 0x1A))
textbox(slide, '© 2026 Gemicates  |  Confidential Product Presentation  |  v2.0',
        0, 6.92, 13.33, 0.4,
        size=10, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/var/www/tracker/docs/Gemicates_Presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
